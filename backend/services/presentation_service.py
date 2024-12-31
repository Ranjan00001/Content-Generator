from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import os, re
import json
import uuid
from configs.config import themes
import logging
from models.generative_model import get_model
from models.generative_model import SlideLayout, generate_slide_prompts

class PresentationService:
    def __init__(self):
        self.model = get_model()
        self.logger = logging.getLogger(__name__)
        logging.basicConfig(level=logging.INFO)
        self.generated_files_dir = "data/generated"
        os.makedirs(self.generated_files_dir, exist_ok=True)

    @staticmethod
    def hex_to_rgbcolor(hex_color):
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def set_slide_background(self, slide, hex_color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.hex_to_rgbcolor(hex_color)

    # def add_text(self, slide, text, font_name, font_size, position, color="000000"):
    #     textbox = slide.shapes.add_textbox(*position)
    #     text_frame = textbox.text_frame
    #     text_frame.word_wrap = True
    #     p = text_frame.add_paragraph()
    #     p.text = text
    #     p.font.name = font_name
    #     p.font.size = Pt(font_size)
    #     p.font.color.rgb = self.hex_to_rgbcolor(color)


    @staticmethod
    def clean_text(text):
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        text = re.sub(r'__(.*?)__', r'\1', text)
        text = re.sub(r'_(.*?)_', r'\1', text)
        return text

    def fetch_content_from_gemini(self, topic, num_slides, layouts):
        """
        Fetches slide content from the Gemini model based on the topic, number of slides, and layouts.

        Args:
            topic (str): Topic for the slides.
            num_slides (int): Number of slides to generate.
            layouts (list[SlideLayout]): List of slide layouts.

        Returns:
            list[dict]: List of content for each slide.
        """
        # Generate the prompt using the helper function
        prompt = generate_slide_prompts(topic, num_slides, layouts)
        # print('prompt', prompt)
        # Generate content using the model
        try:
            response = self.model.generate_content(prompt).to_dict()
        except Exception as e:
            self.logger.error("Error generating content from Gemini: %s", e)
            raise RuntimeError("Content generation failed") from e

        text = response['candidates'][0]['content']['parts'][0]['text'][7:-5]
        slide_data = []
        try:
            # Parse response text as JSON
            slide_data = json.loads(text)
        except json.JSONDecodeError as e:
            self.logger.warning("Failed to parse part text as JSON: %s", text)

        # Validate the generated content
        if not slide_data:
            self.logger.error("No content generated for the topic: %s", topic)
            raise ValueError("Generated content is empty or invalid.")

        return slide_data
    
    def add_title(self, slide, title, theme):
        self.set_slide_background(slide, theme["background_color"])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = theme["font"]
        title_shape.text_frame.paragraphs[0].font.size = Pt(theme["font_size"] + 10)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.hex_to_rgbcolor(theme["title_color"])

    def add_subtitle(self, slide, subtitle, theme):
        slide.placeholders[1].text = subtitle

    def add_bullet_points(self, slide, title, points, theme):
        self.set_slide_background(slide, theme["background_color"])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = theme["font"]
        title_shape.text_frame.paragraphs[0].font.size = Pt(theme["font_size"] + 5)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.hex_to_rgbcolor(theme["title_color"])
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.bullet = True
            p.level = 0
            p.font.name = theme["font"]
            p.font.size = Pt(theme["font_size"])
            p.font.color.rgb = self.hex_to_rgbcolor(theme["content_color"])

    def add_two_column(self, slide, title, left, right, theme):
        self.set_slide_background(slide, theme["background_color"])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = theme["font"]
        title_shape.text_frame.paragraphs[0].font.size = Pt(theme["font_size"] + 5)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.hex_to_rgbcolor(theme["title_color"])
        left_box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(4), Inches(4))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        for point in left:
            p = left_tf.add_paragraph()
            p.text = point
            p.bullet = True
            p.font.name = theme["font"]
            p.font.size = Pt(theme["font_size"])
            p.font.color.rgb = self.hex_to_rgbcolor(theme["content_color"])

        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.7), Inches(4), Inches(4))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        for point in right:
            p = right_tf.add_paragraph()
            p.text = point
            p.bullet = True
            p.font.name = theme["font"]
            p.font.size = Pt(theme["font_size"])
            p.font.color.rgb = self.hex_to_rgbcolor(theme["content_color"])

    def add_content_with_image(self, slide, title, content, image, theme):
        self.set_slide_background(slide, theme["background_color"])
                
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = theme["font"]
        title_shape.text_frame.paragraphs[0].font.size = Pt(theme["font_size"] + 5)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.hex_to_rgbcolor(theme["title_color"])
                
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(5), Inches(4))
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        content_tf.text = content
        content_tf.paragraphs[0].font.name = theme["font"]
        content_tf.paragraphs[0].font.size = Pt(theme["font_size"] + 5)
        content_tf.paragraphs[0].font.color.rgb = self.hex_to_rgbcolor(theme["title_color"])

        image_path = image
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(7), Inches(1.7), width=Inches(2), height=Inches(2))
        else:
            placeholder = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(2), Inches(2))
            placeholder.text = "Image Not Found"
            placeholder.text_frame.paragraphs[0].font.size = Pt(theme["font_size"])
            placeholder.text_frame.paragraphs[0].font.color.rgb = self.hex_to_rgbcolor(theme["content_color"])

    def create_presentation(self, content, theme_name="default"):
        prs = Presentation()
        theme = themes.get(theme_name, themes["default"])

        for slide_content in content:
            layout = slide_content.get("layout", SlideLayout.BULLET_POINTS.value)

            if layout == SlideLayout.TITLE.value:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                self.add_title(slide, slide_content['title'], theme)
                if "subtitle" in slide_content:
                    self.add_subtitle(slide, slide_content['subtitle'], theme)

            elif layout == SlideLayout.BULLET_POINTS.value:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                self.add_bullet_points(slide, slide_content['title'], slide_content['points'], theme)
                
            elif layout == SlideLayout.TWO_COLUMN.value:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide_content.get('title', '')
                left_point = slide_content.get('left_points', [])
                right_points = slide_content.get('right_points', [])
                self.add_two_column(slide, title, left_point, right_points, theme)

            elif layout == SlideLayout.CONTENT_WITH_IMAGE.value:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                self.add_content_with_image(slide, slide_content.get('title'), slide_content.get('content'), slide_content.get('image_path'), theme)

            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                self.set_slide_background(slide, theme["background_color"])
                title_shape = slide.shapes.title
                title_shape.text = slide_content["title"]
                title_shape.text_frame.paragraphs[0].font.name = theme["font"]
                title_shape.text_frame.paragraphs[0].font.size = Pt(theme["font_size"] + 5)
                title_shape.text_frame.paragraphs[0].font.color.rgb = self.hex_to_rgbcolor(theme["title_color"])
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                for point in slide_content.get("points", []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.bullet = True
                    p.level = 0
                    p.font.name = theme["font"]
                    p.font.size = Pt(theme["font_size"])
                    p.font.color.rgb = self.hex_to_rgbcolor(theme["content_color"])

        return prs

    # def save_presentation(self, prs, topic, num_slides, theme, layouts):
    #     presentation_id = str(uuid.uuid4())
    #     pptx_path = os.path.join(self.generated_files_dir, f"{presentation_id}.pptx")
    #     prs.save(pptx_path)

    #     presentation_details = {
    #         "id": presentation_id,
    #         "topic": topic,
    #         "num_slides": num_slides,
    #         "theme": theme,
    #         "layouts": layouts,
    #         "download_url": f"/api/v1/presentations/{presentation_id}/download",
    #     }

    #     json_path = os.path.join(self.generated_files_dir, f"{presentation_id}.json")
    #     with open(json_path, 'w') as json_file:
    #         json.dump(presentation_details, json_file, indent=4)

    #     return presentation_details
